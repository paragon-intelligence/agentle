import subprocess
import os
import tempfile
from collections.abc import MutableSequence
from pathlib import Path
from typing import Literal, cast, override

from rsb.models.field import Field

from agentle.agents.agent import Agent
from agentle.generations.models.message_parts.file import FilePart
from agentle.generations.models.structured_outputs_store.visual_media_description import (
    VisualMediaDescription,
)
from agentle.generations.providers.base.generation_provider import GenerationProvider
from agentle.generations.providers.google.google_genai_generation_provider import (
    GoogleGenaiGenerationProvider,
)
from agentle.parsing.document_parser import DocumentParser
from agentle.parsing.factories.visual_description_agent_factory import (
    visual_description_agent_factory,
)
from agentle.parsing.image import Image
from agentle.parsing.parsed_document import ParsedDocument
from agentle.parsing.parses import parses
from agentle.parsing.section_content import SectionContent


@parses("ppt", "pptx", "pptm")
class PptxFileParser(DocumentParser):
    strategy: Literal["high", "low"] = Field(default="high")

    visual_description_agent: Agent[VisualMediaDescription] = Field(
        default_factory=visual_description_agent_factory,
    )
    """
    The agent to use for generating the visual description of the document.
    Useful when you want to customize the prompt for the visual description.
    """

    multi_modal_provider: GenerationProvider = Field(
        default_factory=GoogleGenaiGenerationProvider,
    )
    """
    The multi-modal provider to use for generating the visual description of the document.
    Useful when you want us to customize the prompt for the visual description.
    """

    @override
    async def parse_async(
        self,
        document_path: str,
    ) -> ParsedDocument:
        import hashlib

        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from pptx.presentation import Presentation as PptxPresentation
        from pptx.shapes.autoshape import Shape
        from pptx.shapes.picture import Picture

        path = Path(document_path)
        converted_pptx_file: str | None = None
        if path.suffix in {".ppt", ".pptm"}:
            converted_pptx_file = self._convert_to_pptx(document_path)

        prs: PptxPresentation = Presentation(converted_pptx_file or document_path)
        sections: MutableSequence[SectionContent] = []
        processed_images: dict[str, tuple[str, str]] = {}

        for slide_index, slide in enumerate(prs.slides, start=1):
            slide_texts: list[str] = []
            slide_images: list[tuple[str, bytes, str]] = []  # (name, data, hash)

            for shape in slide.shapes:
                if shape.has_text_frame:
                    shape_with_text = cast(Shape, shape)
                    text_str: str = shape_with_text.text
                    slide_texts.append(text_str)

                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    picture_shape = cast(Picture, shape)
                    image_blob: bytes = picture_shape.image.blob
                    image_hash = hashlib.sha256(image_blob).hexdigest()
                    image_name: str = (
                        shape.name or f"slide_{slide_index}_img_{image_hash[:8]}"
                    )
                    slide_images.append((image_name, image_blob, image_hash))

            combined_text: str = "\n".join(slide_texts)
            final_images: MutableSequence[Image] = []
            image_descriptions: MutableSequence[str] = []

            if self.visual_description_agent and self.strategy == "high":
                for img_idx, (image_name, image_blob, image_hash) in enumerate(
                    slide_images, start=1
                ):
                    is_cached = image_hash in processed_images
                    if is_cached:
                        cached_md, cached_ocr = processed_images[image_hash]
                        image_descriptions.append(
                            f"Slide {slide_index} - Image {img_idx}: {cached_md}"
                        )
                        final_images.append(
                            Image(
                                name=image_name,
                                contents=image_blob,
                                ocr_text=cached_ocr,
                            )
                        )
                        continue

                    agent_input = FilePart(
                        mime_type=Path(image_name).suffix,
                        data=image_blob,
                    )
                    agent_response = await self.visual_description_agent.run_async(
                        agent_input
                    )
                    image_md: str = agent_response.parsed.md
                    image_ocr = agent_response.parsed.ocr_text

                    processed_images[image_hash] = (image_md, image_ocr or "")
                    image_descriptions.append(
                        f"Slide {slide_index} - Image {img_idx}: {image_md}"
                    )
                    final_images.append(
                        Image(name=image_name, contents=image_blob, ocr_text=image_ocr)
                    )

                if image_descriptions:
                    combined_text += "\n\n" + "\n".join(image_descriptions)

            section_content = SectionContent(
                number=slide_index,
                text=combined_text,
                md=combined_text,
                images=final_images,
            )
            sections.append(section_content)

        return ParsedDocument(
            name=path.name,
            sections=sections,
        )

    def _convert_to_pptx(self, document_path: str) -> str:
        def _is_libreoffice_installed() -> bool:
            try:
                subprocess.run(
                    ["libreoffice", "--version"],
                    stdout=subprocess.DEVNULL,
                    stderr=subprocess.DEVNULL,
                    check=True,
                )
                return True
            except (subprocess.CalledProcessError, FileNotFoundError):
                return False

        if not _is_libreoffice_installed():
            raise RuntimeError("LibreOffice not found in system PATH")

        with tempfile.TemporaryDirectory() as temp_dir:
            # Write input file to temporary directory
            input_filename = Path(document_path).name
            input_path = os.path.join(temp_dir, input_filename)
            with (
                open(document_path, "rb") as src_file,
                open(input_path, "wb") as dst_file,
            ):
                dst_file.write(src_file.read())

            # Run LibreOffice conversion
            try:
                subprocess.run(
                    [
                        "libreoffice",
                        "--headless",
                        "--convert-to",
                        "pptx",
                        "--outdir",
                        temp_dir,
                        input_path,
                    ],
                    check=True,
                    stdout=subprocess.DEVNULL,
                    stderr=subprocess.PIPE,
                    timeout=60,
                )
            except subprocess.CalledProcessError as e:
                error_msg = e.stderr.decode().strip() if e.stderr else "Unknown error"
                raise RuntimeError(f"Conversion failed: {error_msg}") from e
            except subprocess.TimeoutExpired:
                raise RuntimeError("Conversion timed out after 60 seconds")

            # Determine output file path
            output_filename = Path(input_filename).stem + ".pptx"
            output_path = os.path.join(temp_dir, output_filename)

            if not os.path.exists(output_path):
                available_files = os.listdir(temp_dir)
                raise RuntimeError(
                    f"Converted file not found at {output_path}. Found files: {available_files}"
                )

            # Read the converted file and create a temporary file in the system temp directory
            with open(output_path, "rb") as f:
                content = f.read()

            temp_output = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
            temp_output.write(content)
            temp_output.close()

            return temp_output.name
