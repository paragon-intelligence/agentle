from agentle.parsing.parses import parses


@_parses("doc", "docx")
class DocxFileParser(FileParser, frozen=True, tag="docx"):
    @ensure_module_installed("docx", "intellibricks[files]")
    @override
    async def parse_async(
        self,
        file: RawFile,
    ) -> ParsedFile:
        import hashlib

        from docx import Document

        with tempfile.TemporaryDirectory() as temp_dir:
            file_path = f"{temp_dir}/{file.name}"
            file.save_to_file(file_path)

            document = Document(file_path)
            image_cache: dict[str, tuple[str, str]] = {}  # (md, ocr_text)

            paragraph_texts = [p.text for p in document.paragraphs if p.text.strip()]
            doc_text = "\n".join(paragraph_texts)

            doc_images: list[tuple[str, bytes]] = []
            for rel in document.part._rels.values():  # type: ignore[reportPrivateUsage]
                if "image" in rel.reltype:
                    image_part = rel.target_part
                    image_name = image_part.partname.split("/")[-1]
                    image_bytes = image_part.blob
                    doc_images.append((image_name, image_bytes))

            final_images: list[Image] = []
            image_descriptions: list[str] = []
            if self.visual_description_agent and self.strategy == "high":
                for idx, (image_name, image_bytes) in enumerate(doc_images, start=1):
                    image_hash = hashlib.sha256(image_bytes).hexdigest()

                    if image_hash in image_cache:
                        cached_md, cached_ocr = image_cache[image_hash]
                        image_md = cached_md
                        ocr_text = cached_ocr
                    else:
                        agent_input = ImageFilePart(
                            mime_type=bytes_to_mime(image_bytes),
                            data=image_bytes,
                        )
                        agent_response = await self.visual_description_agent.run_async(
                            agent_input
                        )
                        image_md = agent_response.parsed.final_answer.md
                        ocr_text = agent_response.parsed.final_answer.ocr_text
                        image_cache[image_hash] = (image_md, ocr_text)

                    image_descriptions.append(f"Docx Image {idx}: {image_md}")
                    final_images.append(
                        Image(
                            name=image_name,
                            contents=image_bytes,
                            ocr_text=ocr_text,
                        )
                    )

                if image_descriptions:
                    doc_text += "\n\n" + "\n".join(image_descriptions)

            return ParsedFile(
                name=file.name,
                sections=[
                    SectionContent(
                        number=1,
                        text=doc_text,
                        md=doc_text,
                        images=final_images,
                    )
                ],
            )


@_parses("ppt", "pptx", "pptm")
class PptxFileParser(FileParser, frozen=True, tag="pptx"):
    @ensure_module_installed("pptx", "intellibricks[files]")
    @override
    async def parse_async(
        self,
        file: RawFile,
    ) -> ParsedFile:
        import hashlib

        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from pptx.presentation import Presentation as PptxPresentation
        from pptx.shapes.autoshape import Shape
        from pptx.shapes.picture import Picture

        with tempfile.TemporaryDirectory() as temp_dir:
            file_path = f"{temp_dir}/{file.name}"
            if file.extension in {"ppt", "pptm"}:
                converted_pptx_file = self._convert_to_pptx(file)
                converted_pptx_file.save_to_file(file_path)
            else:
                file.save_to_file(file_path)

            prs: PptxPresentation = Presentation(file_path)
            sections: list[SectionContent] = []
            processed_images: dict[str, tuple[str, str]] = {}

            for slide_index, slide in enumerate(prs.slides, start=1):
                slide_texts: list[str] = []
                slide_images: list[tuple[str, bytes, str]] = []  # (name, data, hash)

                for shape in slide.shapes:
                    if shape.has_text_frame:
                        shape_with_text = cast(Shape, shape)
                        text_str: str = shape_with_text.text
                        slide_texts.append(text_str)

                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        picture_shape = cast(Picture, shape)
                        image_blob: bytes = picture_shape.image.blob
                        image_hash = hashlib.sha256(image_blob).hexdigest()
                        image_name: str = (
                            shape.name or f"slide_{slide_index}_img_{image_hash[:8]}"
                        )
                        slide_images.append((image_name, image_blob, image_hash))

                combined_text: str = "\n".join(slide_texts)
                final_images: list[Image] = []
                image_descriptions: list[str] = []

                if self.visual_description_agent and self.strategy == "high":
                    for img_idx, (image_name, image_blob, image_hash) in enumerate(
                        slide_images, start=1
                    ):
                        is_cached = image_hash in processed_images
                        if is_cached:
                            cached_md, cached_ocr = processed_images[image_hash]
                            image_descriptions.append(
                                f"Slide {slide_index} - Image {img_idx}: {cached_md}"
                            )
                            final_images.append(
                                Image(
                                    name=image_name,
                                    contents=image_blob,
                                    ocr_text=cached_ocr,
                                )
                            )
                            continue

                        agent_input = ImageFilePart(
                            mime_type=bytes_to_mime(image_blob),
                            data=image_blob,
                        )
                        agent_response = await self.visual_description_agent.run_async(
                            agent_input
                        )
                        image_md: str = agent_response.parsed.final_answer.md
                        image_ocr = agent_response.parsed.final_answer.ocr_text

                        processed_images[image_hash] = (image_md, image_ocr)
                        image_descriptions.append(
                            f"Slide {slide_index} - Image {img_idx}: {image_md}"
                        )
                        final_images.append(
                            Image(
                                name=image_name, contents=image_blob, ocr_text=image_ocr
                            )
                        )

                    if image_descriptions:
                        combined_text += "\n\n" + "\n".join(image_descriptions)

                section_content = SectionContent(
                    number=slide_index,
                    text=combined_text,
                    md=combined_text,
                    images=final_images,
                )
                sections.append(section_content)

            return ParsedFile(
                name=file.name,
                sections=sections,
            )

    def _convert_to_pptx(self, file: RawFile) -> RawFile:
        """Convert PowerPoint files (.ppt/.pptm) to .pptx format and return as RawFile.

        Args:
            file: RawFile instance containing the input file data.

        Returns:
            RawFile instance containing converted content.

        Raises:
            RuntimeError: If conversion fails or LibreOffice not installed.
        """

        def _is_libreoffice_installed() -> bool:
            try:
                subprocess.run(
                    ["libreoffice", "--version"],
                    stdout=subprocess.DEVNULL,
                    stderr=subprocess.DEVNULL,
                    check=True,
                )
                return True
            except (subprocess.CalledProcessError, FileNotFoundError):
                return False

        if not _is_libreoffice_installed():
            raise RuntimeError("LibreOffice not found in system PATH")

        with tempfile.TemporaryDirectory() as temp_dir:
            # Write input file to temporary directory
            input_path = os.path.join(temp_dir, file.name)
            with open(input_path, "wb") as f:
                f.write(file.contents)

            # Run LibreOffice conversion
            try:
                subprocess.run(
                    [
                        "libreoffice",
                        "--headless",
                        "--convert-to",
                        "pptx",
                        "--outdir",
                        temp_dir,
                        input_path,
                    ],
                    check=True,
                    stdout=subprocess.DEVNULL,
                    stderr=subprocess.PIPE,
                    timeout=60,
                )
            except subprocess.CalledProcessError as e:
                error_msg = e.stderr.decode().strip() if e.stderr else "Unknown error"
                raise RuntimeError(f"Conversion failed: {error_msg}") from e
            except subprocess.TimeoutExpired:
                raise RuntimeError("Conversion timed out after 60 seconds")

            # Determine output file path
            output_filename = Path(file.name).stem + ".pptx"
            output_path = os.path.join(temp_dir, output_filename)

            if not os.path.exists(output_path):
                available_files = os.listdir(temp_dir)
                raise RuntimeError(
                    f"Converted file not found at {output_path}. Found files: {available_files}"
                )

            # Read converted file and return as RawFile
            return RawFile.from_file_path(output_path)


@_parses("xls", "xlsx")
class ExcelFileParser(FileParser, frozen=True, tag="excel"):
    @ensure_module_installed("openpyxl", "intellibricks[files]")
    @override
    async def parse_async(
        self,
        file: RawFile,
    ) -> ParsedFile:
        import csv
        import io

        from openpyxl import Workbook, load_workbook

        with tempfile.TemporaryDirectory() as temp_dir:
            file_path = f"{temp_dir}/{file.name}"
            file.save_to_file(file_path)

            wb: Workbook = load_workbook(file_path, data_only=True)
            sections: list[SectionContent] = []

            for sheet_index, sheet in enumerate(wb.worksheets, start=1):
                # Gather structured data
                rows: list[list[str]] = []
                row_texts: list[str] = []
                for row in sheet.iter_rows(values_only=True):
                    # Process cell values
                    cell_values = [
                        str(cell) if cell is not None else "" for cell in row
                    ]
                    rows.append(cell_values)
                    row_texts.append("\t".join(cell_values))

                combined_text = "\n".join(row_texts)

                # Generate CSV content
                csv_buffer = io.StringIO()
                csv_writer = csv.writer(csv_buffer)
                csv_writer.writerows(rows)
                csv_str = csv_buffer.getvalue().strip()

                # Process images
                sheet_images: list[tuple[str, bytes]] = []
                if hasattr(sheet, "_images"):
                    image_list = getattr(sheet, "_images", [])
                    for img_idx, img in enumerate(image_list, start=1):
                        img_data = getattr(img, "_data", None)
                        if img_data is not None:
                            image_name = f"{sheet.title}_img_{img_idx}.png"
                            sheet_images.append((image_name, img_data))

                final_images: list[Image] = []
                # Generate image descriptions if needed
                if self.visual_description_agent and self.strategy == "high":
                    image_descriptions: list[str] = []
                    for img_idx, image_obj in enumerate(sheet_images, start=1):
                        agent_input = ImageFilePart(
                            mime_type=bytes_to_mime(image_obj[1]),
                            data=image_obj[1],
                        )
                        agent_response = await self.visual_description_agent.run_async(
                            agent_input
                        )
                        image_md = agent_response.parsed.final_answer.md
                        image_descriptions.append(
                            f"Worksheet {sheet.title} - Image {img_idx}: {image_md}"
                        )
                        final_images.append(
                            Image(
                                name=image_obj[0],
                                contents=image_obj[1],
                                ocr_text=agent_response.parsed.final_answer.ocr_text,
                            )
                        )

                    if image_descriptions:
                        combined_text += "\n\n" + "\n".join(image_descriptions)

                # Create table page item
                table_item = TablePageItem(
                    md=combined_text, rows=rows, csv=csv_str, is_perfect_table=True
                )

                section_content = SectionContent(
                    number=sheet_index,
                    text=combined_text,
                    md=combined_text,
                    images=final_images,
                    items=[table_item],
                )
                sections.append(section_content)

            return ParsedFile(
                name=file.name,
                sections=sections,
            )


@_parses("txt", "alg")
class TxtFileParser(FileParser, frozen=True, tag="txt"):
    @override
    async def parse_async(self, file: RawFile) -> ParsedFile:
        text_content = file.contents.decode("utf-8", errors="replace")

        page_content = SectionContent(
            number=1,
            text=text_content,
            md=text_content,
        )

        return ParsedFile(
            name=file.name,
            sections=[page_content],
        )


@_parses("gif")
class AnimatedImageFileParser(FileParser, frozen=True, tag="animated_image"):
    @override
    async def parse_async(
        self,
        file: RawFile,
    ) -> ParsedFile:
        from PIL import Image as PILImage

        with tempfile.TemporaryDirectory() as temp_dir:
            file_path = f"{temp_dir}/{file.name}"
            file.save_to_file(file_path)

            # Safety check: only proceed if it's a .gif
            # or you can attempt detection based on file headers
            extension = file.extension
            if extension not in {"gif"}:
                raise ValueError("AnimatedImageFileParser only supports .gif files.")

            # --- 1. Load all frames from the GIF ---
            frames: list[PILImage.Image] = []
            with PILImage.open(file_path) as gif_img:
                try:
                    while True:
                        frames.append(gif_img.copy())
                        gif_img.seek(gif_img.tell() + 1)
                except EOFError:
                    pass  # we've reached the end of the animation

            num_frames = len(frames)
            if num_frames == 0:
                # No frames => no content
                return ParsedFile(name=file.name, sections=[])

            # --- 2. Pick up to 3 frames, splitting the GIF into 3 segments ---
            # If there are fewer than 3 frames, just use them all.
            # If more than 3, pick three frames spaced across the animation.

            if num_frames <= 3:
                selected_frames = frames
            else:
                # Example approach: pick near 1/3, 2/3, end
                idx1 = max(0, (num_frames // 3) - 1)
                idx2 = max(0, (2 * num_frames // 3) - 1)
                idx3 = num_frames - 1
                # Ensure distinct indexes
                unique_indexes = sorted(set([idx1, idx2, idx3]))
                selected_frames = [frames[i] for i in unique_indexes]

            # --- 3. Convert each selected frame to PNG and (optionally) describe it ---
            pages: list[SectionContent] = []
            for i, frame in enumerate(selected_frames, start=1):
                # Convert frame to PNG in-memory
                png_buffer = io.BytesIO()
                # Convert to RGBA if needed
                if frame.mode not in ("RGB", "RGBA"):
                    frame = frame.convert("RGBA")
                frame.save(png_buffer, format="PNG")
                png_bytes = png_buffer.getvalue()

                frame_image_ocr: str | None = None
                # If strategy is HIGH, pass the frame to the agent
                text_description = ""
                if self.visual_description_agent and self.strategy == "high":
                    agent_input = ImageFilePart(
                        mime_type=bytes_to_mime(png_bytes),
                        data=png_bytes,
                    )
                    agent_response = await self.visual_description_agent.run_async(
                        agent_input
                    )
                    frame_image_ocr = agent_response.parsed.final_answer.ocr_text
                    text_description = agent_response.parsed.final_answer.md

                # Create an Image object
                frame_image = Image(
                    name=f"{file.name}-frame{i}.png",
                    contents=png_bytes,
                    ocr_text=frame_image_ocr,
                )
                # Each frame is its own "page" in the final doc
                page_content = SectionContent(
                    number=i,
                    text=text_description,
                    md=text_description,
                    images=[frame_image],
                )
                pages.append(page_content)

            # --- 4. Return the multi-page ParsedFile ---
            return ParsedFile(
                name=file.name,
                sections=pages,
            )


@_parses("flac", "mp3", "mpeg", "mpga", "m4a", "ogg", "wav", "webm")
class AudioFileParser(FileParser, frozen=True, tag="audio"):
    async def parse_async(self, file: RawFile) -> ParsedFile:
        if self.audio_description_agent is None:
            raise ValueError("No audio description agent provided.")

        file_contents: bytes = file.contents
        file_extension = file.extension

        if file_extension in {
            "flac",
            "mpeg",
            "mpga",
            "m4a",
            "ogg",
            "wav",
            "webm",
        }:
            import aiofiles.os as aios
            from aiofiles import open as aio_open

            self._check_ffmpeg_installed()

            # Generate unique temporary filenames
            input_temp = os.path.join(
                tempfile.gettempdir(),
                f"input_{os.urandom(8).hex()}.{file_extension}",
            )
            output_temp = os.path.join(
                tempfile.gettempdir(), f"output_{os.urandom(8).hex()}.mp3"
            )

            # Write input file asynchronously
            async with aio_open(input_temp, "wb") as f:
                await f.write(file_contents)

            # Build FFmpeg command
            command = [
                "ffmpeg",
                "-hide_banner",
                "-loglevel",
                "error",  # Suppress unnecessary logs
                "-y",  # Overwrite output file if exists
                "-i",
                input_temp,
                "-codec:a",
                "libmp3lame",
                "-q:a",
                "2",  # Quality preset (0-9, 0=best)
                output_temp,
            ]

            # Execute FFmpeg
            process = await asyncio.create_subprocess_exec(
                *command, stdout=asyncio.subprocess.PIPE, stderr=asyncio.subprocess.PIPE
            )

            _, stderr = await process.communicate()

            # Handle conversion errors
            if process.returncode != 0:
                await aios.remove(input_temp)
                if await aios.path.exists(output_temp):
                    await aios.remove(output_temp)
                raise RuntimeError(
                    f"Audio conversion failed: {stderr.decode().strip()}"
                )

            # Read converted file
            async with aio_open(output_temp, "rb") as f:
                file_contents = await f.read()

            # Cleanup temporary files
            await aios.remove(input_temp)
            await aios.remove(output_temp)

        transcription = self.audio_description_agent.run(
            AudioFilePart(data=file_contents, mime_type=bytes_to_mime(file_contents))
        )

        return ParsedFile(
            name=file.name,
            sections=[
                SectionContent(
                    number=1,
                    text=transcription.audio_transcription.text
                    if transcription.audio_transcription is not None
                    else self._could_not_transcript(),
                    md=transcription.parsed.final_answer.md,
                    images=[],
                )
            ],
        )

    def _could_not_transcript(self) -> Never:
        raise ValueError("Could not transcribe the audio")

    def _check_ffmpeg_installed(self) -> None:
        import subprocess

        try:
            result = subprocess.run(
                ["ffmpeg", "-version"],
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                text=True,
            )
            exception_logger.exception("FFmpeg is not installed or not in PATH.")
            if result.returncode != 0:
                raise RuntimeError()
        except FileNotFoundError:
            exception_logger.exception("FFmpeg is not installed or not in PATH.")
            raise RuntimeError()


@parses("mp4")
class VideoFileParser(FileParser, frozen=True, tag="video"):
    async def parse_async(self, file: RawFile) -> ParsedFile:
        if self.visual_description_agent is None:
            raise ValueError("No visual description agent provided.")

        extension = file.extension
        if extension != "mp4":
            raise ValueError("VideoFileParser only supports .mp4 files.")

        file_contents = file.contents
        visual_media_description = await self.visual_description_agent.run_async(
            VideoFilePart(data=file_contents, mime_type=bytes_to_mime(file_contents))
        )

        return ParsedFile(
            name=file.name,
            sections=[
                SectionContent(
                    number=1,
                    text=visual_media_description.parsed.final_answer.md,
                    md=visual_media_description.parsed.final_answer.md,
                    images=[],
                )
            ],
        )
